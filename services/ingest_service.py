"""CorporateBrain – document ingestion service.

Pipeline: download / receive bytes  →  extract text  →  clean  →  smart chunk
          →  embed  →  store in Supabase.

Supported file types:
  .docx  — paragraphs + tables + footnotes + text boxes
  .pdf   — pdfplumber (layout-aware) with pytesseract OCR fallback for scans
  .xlsx  — every sheet as labelled rows
  .pptx  — slide text + speaker notes
  .txt / .md — plain UTF-8

Key design rules:
  - Markdown tables kept intact as atomic chunks (never split mid-row).
  - Document title / filename prepended to every chunk for cross-doc synthesis.
  - Chunk size 1500 chars so granular facts land in their own retrievable chunk.
  - Overlap 200 chars catches facts that straddle chunk boundaries.
  - Text cleaned before chunking: artifacts, watermarks, rogue page numbers removed.
"""

from __future__ import annotations

import io
import logging
import os
import re
from typing import List

import httpx
from docx import Document as DocxDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import AsyncOpenAI

from services.db import supabase

logger = logging.getLogger(__name__)

# ── OpenAI-compatible client ─────────────────────────────────────────────────
_openai = AsyncOpenAI(
    base_url="https://models.inference.ai.azure.com",
    api_key=os.environ.get("GITHUB_TOKEN", ""),
)

EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIMENSIONS = 1536

# Repeated watermark words that pollute chunks if left in
_WATERMARK_WORDS = {"CONFIDENTIAL", "DRAFT", "PRIVILEGED", "SAMPLE", "COPY"}


# ── Text cleaning ─────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    """Normalise extracted text before chunking.

    Removes common extraction artifacts:
      - Tab characters collapsed to single space
      - Multiple consecutive spaces collapsed
      - Isolated page numbers (lone digits on their own line)
      - Repeated watermark words (e.g. CONFIDENTIAL CONFIDENTIAL CONFIDENTIAL)
      - Windows-style line endings normalised
      - Excessive blank lines (3+ collapsed to 2)
    """
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\t+", " ", text)
    text = re.sub(r"[ ]{2,}", " ", text)
    text = re.sub(r"(?m)^\s*\d{1,4}\s*$", "", text)
    for word in _WATERMARK_WORDS:
        text = re.sub(rf"(\b{word}\b\s*){{2,}}", f"{word} ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ── Text extraction ───────────────────────────────────────────────────────────

async def _download_file(url: str) -> bytes:
    async with httpx.AsyncClient(timeout=60.0) as client:
        response = await client.get(url)
        response.raise_for_status()
        return response.content


def _extract_text_docx(docx_bytes: bytes) -> str:
    """Extract text from .docx: paragraphs + tables + footnotes + text boxes.

    Walks the document body XML in document order so content appears in the
    correct reading sequence. Also surfaces footnote text (common in legal
    documents) and text boxes which python-docx skips by default.
    """
    from docx.oxml.ns import qn

    doc = DocxDocument(io.BytesIO(docx_bytes))
    parts: List[str] = []

    # ── Body: paragraphs and tables in document order ────────────────────────
    for block in doc.element.body:
        tag = block.tag

        if tag == qn("w:p"):
            text = "".join(
                node.text or ""
                for node in block.iter()
                if node.tag == qn("w:t")
            )
            if text.strip():
                parts.append(text.strip())

        elif tag == qn("w:tbl"):
            # Emit each row as pipe-delimited markdown so the table splitter
            # keeps it atomic.
            for row in block.findall(".//" + qn("w:tr")):
                cells: List[str] = []
                for cell in row.findall(".//" + qn("w:tc")):
                    cell_text = "".join(
                        node.text or ""
                        for node in cell.iter()
                        if node.tag == qn("w:t")
                    ).strip()
                    cells.append(cell_text)
                if any(cells):
                    parts.append("| " + " | ".join(cells) + " |")

    # ── Footnotes ────────────────────────────────────────────────────────────
    try:
        footnotes_part = doc.part.footnotes
        if footnotes_part is not None:
            fn_root = footnotes_part._element
            for fn in fn_root.findall(qn("w:footnote")):
                fn_id = fn.get(qn("w:id"), "")
                if fn_id in ("-1", "0"):
                    continue
                fn_text = "".join(
                    node.text or ""
                    for node in fn.iter()
                    if node.tag == qn("w:t")
                ).strip()
                if fn_text:
                    parts.append(f"[Footnote] {fn_text}")
    except Exception as exc:
        logger.debug("Footnote extraction skipped: %s", exc)

    return _clean_text("\n\n".join(parts))


def _extract_text_plain(raw: bytes) -> str:
    return _clean_text(raw.decode("utf-8", errors="replace"))


def _extract_text_pdf(pdf_bytes: bytes) -> str:
    """Extract text from a PDF.

    Strategy:
      1. pdfplumber — layout-aware, preserves table structure and handles
         multi-column documents better than pypdf.
      2. If pdfplumber returns fewer than 100 chars total the PDF is likely a
         scanned image. Fall back to pytesseract OCR via pdf2image.
      3. If OCR dependencies are unavailable a clear warning is logged and an
         empty string returned so the caller surfaces a friendly error.
    """
    try:
        import pdfplumber
    except ImportError:
        raise ImportError(
            "pdfplumber is required for PDF extraction. "
            "Add 'pdfplumber' to requirements.txt."
        )

    pages: List[str] = []

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages:
            page_parts: List[str] = []

            # Tables first — extract as markdown rows
            for table in page.extract_tables():
                for row in table:
                    cells = [str(cell or "").strip() for cell in row]
                    if any(cells):
                        page_parts.append("| " + " | ".join(cells) + " |")

            # Remaining prose text
            page_text = page.extract_text() or ""
            if page_text.strip():
                page_parts.append(page_text.strip())

            if page_parts:
                pages.append("\n".join(page_parts))

    full_text = "\n\n".join(pages)

    # OCR fallback for scanned PDFs
    if len(full_text.strip()) < 100:
        logger.info("PDF text layer sparse — attempting OCR fallback")
        full_text = _ocr_pdf(pdf_bytes)

    return _clean_text(full_text)


def _ocr_pdf(pdf_bytes: bytes) -> str:
    """Rasterise PDF pages and run Tesseract OCR on each.

    Requires: pdf2image (wraps poppler) and pytesseract + system Tesseract.
    Returns empty string gracefully if either dependency is missing.

    FIX 3: DPI raised from 200 → 300. Agency brand guides, pitch decks, and
    campaign reports frequently use small-point type and dense layouts. 200 DPI
    produces unreliable OCR on anything below ~11pt. 300 DPI is the accepted
    minimum for production-quality text recognition on mixed-content documents.
    """
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError:
        logger.warning(
            "OCR dependencies not installed (pdf2image / pytesseract). "
            "Scanned PDF will produce no text. "
            "Install with: pip install pdf2image pytesseract"
        )
        return ""

    try:
        images = convert_from_bytes(pdf_bytes, dpi=300)  # was 200
        pages = [pytesseract.image_to_string(img) for img in images]
        return "\n\n".join(p.strip() for p in pages if p.strip())
    except Exception as exc:
        logger.warning("OCR failed: %s", exc)
        return ""


def _extract_text_xlsx(xlsx_bytes: bytes) -> str:
    """Extract text from an Excel workbook.

    Each sheet is emitted as a labelled markdown table so the chunker keeps
    rows intact. Empty sheets are skipped.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "openpyxl is required for Excel extraction. "
            "Add 'openpyxl' to requirements.txt."
        )

    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
    parts: List[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        non_empty = [r for r in rows if any(c is not None for c in r)]
        if not non_empty:
            continue

        parts.append(f"## Sheet: {sheet_name}")
        for row in non_empty:
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                parts.append("| " + " | ".join(cells) + " |")

    return _clean_text("\n\n".join(parts))


def _extract_text_pptx(pptx_bytes: bytes) -> str:
    """Extract text from a PowerPoint presentation.

    Each slide is emitted with its slide number, all text shapes, and
    speaker notes so that slide context is preserved in retrieval.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for PowerPoint extraction. "
            "Add 'python-pptx' to requirements.txt."
        )

    prs = Presentation(io.BytesIO(pptx_bytes))
    parts: List[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = [f"## Slide {i}"]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_parts.append(text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""
            if notes_text:
                slide_parts.append(f"[Notes] {notes_text}")

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return _clean_text("\n\n".join(parts))


# ── Document metadata inference ───────────────────────────────────────────────

# FIX 1: Agency-tuned document type classifier.
# Previous version only recognised legal document types (NDA, Engagement Letter,
# Legal Memo, Invoice, Scope of Work). Agency clients primarily upload operational
# documents — briefs, reports, brand guides, strategy docs, case studies — none
# of which matched, causing all chunks to carry an empty doc_type label and losing
# the context benefit of typed metadata in retrieval.

_DOC_TYPE_PATTERNS: List[tuple[str, List[str]]] = [
    # Agency operational documents — checked first as they're the most common
    ("Campaign Brief",      ["campaign brief", "creative brief", "ad brief", "project brief", "campaign overview"]),
    ("Monthly Report",      ["monthly report", "performance report", "campaign report", "monthly performance", "results report"]),
    ("Brand Guidelines",    ["brand guide", "brand guidelines", "tone of voice", "brand standards", "visual identity", "brand identity"]),
    ("Strategy Document",   ["strategy document", "marketing strategy", "growth strategy", "go-to-market", "gtm strategy", "quarterly plan", "90-day plan"]),
    ("Case Study",          ["case study", "client success", "client results", "success story", "client outcome"]),
    ("Pitch Deck",          ["pitch deck", "agency deck", "capabilities deck", "new business deck", "proposal deck"]),
    ("Onboarding Guide",    ["onboarding guide", "onboarding checklist", "client onboarding", "welcome guide", "getting started"]),
    ("SOP / Playbook",      ["standard operating procedure", "sop", "playbook", "process guide", "operating procedure", "how we work"]),
    ("Proposal",            ["proposal", "statement of work", "scope of work", "project proposal", "service proposal"]),
    ("Email Sequence",      ["email sequence", "drip sequence", "nurture sequence", "email campaign", "outreach sequence"]),
    ("Analytics Report",    ["analytics report", "data report", "kpi report", "dashboard report", "metrics report", "performance dashboard"]),
    ("Competitor Analysis", ["competitor analysis", "competitive analysis", "competitive landscape", "market analysis", "competitor review"]),
    # Financial / admin documents
    ("Invoice",             ["invoice", "billing statement", "bill to", "payment due", "amount due"]),
    ("Contract",            ["agreement", "contract", "terms and conditions", "service agreement", "master services agreement"]),
    ("NDA",                 ["non-disclosure", "nda", "confidentiality agreement", "confidential information"]),
    ("Scope of Work",       ["scope of work", "statement of work", "sow", "project scope", "deliverables"]),
]

_PARTY_PATTERNS = [
    # Formal entity names with common suffixes
    re.compile(r'\b([A-Z][A-Za-z &,\.]+?(?:LLC|LLP|Inc|Corp|Ltd|LP))\b'),
    # "between X and Y" constructions
    re.compile(
        r'\bbetween\s+([A-Z][A-Za-z &,\.]+?(?:LLC|LLP|Inc|Corp|Ltd))\s+and\s+'
        r'([A-Z][A-Za-z &,\.]+?(?:LLC|LLP|Inc|Corp|Ltd))',
    ),
]


def _infer_doc_metadata(filename: str, text: str) -> dict:
    """Infer document type and named parties from filename + text sample.

    Checks the filename first (highest signal), then the first 600 characters
    of extracted text. Returns {"document_type": str, "parties": list[str]}.
    """
    sample = (filename + " " + text[:600]).lower()

    doc_type = ""
    for type_label, keywords in _DOC_TYPE_PATTERNS:
        if any(kw in sample for kw in keywords):
            doc_type = type_label
            break

    # Party extraction — works on the raw (mixed-case) text for proper nouns
    raw_sample = filename + " " + text[:600]
    parties: List[str] = []

    for pattern in _PARTY_PATTERNS:
        for match in pattern.finditer(raw_sample):
            for group in match.groups():
                if group:
                    g = group.strip(" ,.")
                    if g and g not in parties:
                        parties.append(g)

    return {"document_type": doc_type, "parties": parties[:4]}


def _split_preserving_tables(text: str) -> List[str]:
    """Split text at paragraph boundaries while keeping markdown table blocks intact.

    A table block is any consecutive sequence of lines that start with '|'.
    These are emitted as single unsplit segments so the chunker never splits
    a table mid-row.
    """
    segments: List[str] = []
    current_table: List[str] = []
    current_prose: List[str] = []

    for line in text.split("\n"):
        if line.strip().startswith("|"):
            if current_prose:
                segments.append("\n".join(current_prose))
                current_prose = []
            current_table.append(line)
        else:
            if current_table:
                segments.append("\n".join(current_table))
                current_table = []
            current_prose.append(line)

    if current_table:
        segments.append("\n".join(current_table))
    if current_prose:
        segments.append("\n".join(current_prose))

    return [s for s in segments if s.strip()]


def _chunk_text(text: str, filename: str = "", doc_metadata: dict | None = None) -> List[str]:
    """Split text into retrieval-optimised chunks.

    Each chunk is prefixed with a rich document label that includes:
      - filename (for citation chips)
      - document_type (Campaign Brief, Monthly Report, Brand Guidelines, etc.)
      - parties (named parties to the document, where present)

    This metadata in the chunk label lets the LLM enforce source isolation —
    knowing that a clause belongs to a specific document type and client context,
    not to a different client's brief or report.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
        chunk_overlap=200,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    segments = _split_preserving_tables(text)

    raw_chunks: List[str] = []
    for segment in segments:
        if "|" in segment and "\n" in segment and len(segment) <= 2000:
            raw_chunks.append(segment)
        else:
            raw_chunks.extend(splitter.split_text(segment))

    doc_meta = doc_metadata or {}
    doc_type = doc_meta.get("document_type", "")
    parties = doc_meta.get("parties", [])

    # Build a rich label so the LLM always knows exactly which document and
    # which context a chunk belongs to.
    label_parts = [f"[Document: {filename}]" if filename else ""]
    if doc_type:
        label_parts.append(f"[Type: {doc_type}]")
    if parties:
        label_parts.append(f"[Parties: {', '.join(parties)}]")
    doc_label = " ".join(p for p in label_parts if p) + "\n\n" if any(label_parts) else ""

    return [doc_label + chunk for chunk in raw_chunks if chunk.strip()]


# ── MIME type map and routing ─────────────────────────────────────────────────

_MIME_MAP = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pdf":  "application/pdf",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt":  "text/plain",
    ".md":   "text/markdown",
}

_SUPPORTED_EXTENSIONS = set(_MIME_MAP.keys())


def _extract_text(ext: str, file_bytes: bytes) -> str:
    """Route to the correct extractor based on file extension."""
    if ext == ".docx":
        return _extract_text_docx(file_bytes)
    elif ext == ".pdf":
        return _extract_text_pdf(file_bytes)
    elif ext == ".xlsx":
        return _extract_text_xlsx(file_bytes)
    elif ext == ".pptx":
        return _extract_text_pptx(file_bytes)
    elif ext in (".txt", ".md"):
        return _extract_text_plain(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}"
        )


# ── Embedding ─────────────────────────────────────────────────────────────────

async def _embed_texts(texts: List[str]) -> List[List[float]]:
    response = await _openai.embeddings.create(
        input=texts,
        model=EMBEDDING_MODEL,
    )
    return [item.embedding for item in response.data]


# ── Main pipeline: URL-based ingest ──────────────────────────────────────────

async def ingest_document(
    document_id: str,
    file_url: str,
    user_id: str | None = None,
) -> int:
    """Run the full ingest pipeline from a URL. Returns chunks stored."""

    filename = ""
    if not user_id:
        doc_row = (
            supabase.table("documents")
            .select("user_id, title")
            .eq("id", document_id)
            .single()
            .execute()
        )
        user_id = doc_row.data.get("user_id")
        filename = doc_row.data.get("title", "")
        if not user_id:
            raise ValueError(
                f"No user_id supplied and document {document_id} has no user_id."
            )

    logger.info("Downloading document %s from %s", document_id, file_url)
    raw_bytes = await _download_file(file_url)

    ext = os.path.splitext(file_url)[1].lower() or ".docx"
    text = _extract_text(ext, raw_bytes)

    if not text.strip():
        raise ValueError(f"Document {document_id} contains no extractable text.")

    chunks = _chunk_text(text, filename=filename, doc_metadata=_infer_doc_metadata(filename, text))
    logger.info("Document %s split into %d chunks", document_id, len(chunks))

    embeddings = await _embed_texts(chunks)

    rows = []
    for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        meta = _infer_doc_metadata(filename, text)
        row = {
            "document_id": document_id,
            "content": chunk,
            "embedding": embedding,
            "chunk_index": idx,
            "metadata": {
                "filename": filename,
                "chunk_index": idx,
                "document_type": meta.get("document_type", ""),
                "parties": meta.get("parties", []),
            },
        }
        if user_id:
            row["user_id"] = user_id
        rows.append(row)

    supabase.table("chunks").insert(rows).execute()
    supabase.table("documents").update({"status": "completed"}).eq("id", document_id).execute()

    logger.info("Document %s ingested successfully (%d chunks)", document_id, len(chunks))
    return len(chunks)


# ── Main pipeline: bytes-based ingest (upload endpoint) ──────────────────────

async def ingest_document_from_bytes(
    document_id: str,
    filename: str,
    file_bytes: bytes,
    user_id: str | None = None,
) -> int:
    """Ingest from raw file bytes (the /upload endpoint). Returns chunks stored."""
    ext = os.path.splitext(filename)[1].lower()

    text = _extract_text(ext, file_bytes)

    if not text.strip():
        raise ValueError(f"Document '{filename}' contains no extractable text.")

    doc_row: dict = {
        "id": document_id,
        "title": filename,
        "status": "processing",
        "mime_type": _MIME_MAP.get(ext, "application/octet-stream"),
        "file_size": len(file_bytes),
    }
    if user_id:
        doc_row["user_id"] = user_id

    try:
        supabase.table("documents").upsert(doc_row).execute()
    except Exception as e:
        logger.warning("Could not upsert document row for %s: %s – continuing", document_id, e)

    doc_meta = _infer_doc_metadata(filename, text)
    chunks = _chunk_text(text, filename=filename, doc_metadata=doc_meta)
    logger.info("Upload '%s' (%s) split into %d chunks", filename, document_id, len(chunks))

    embeddings = await _embed_texts(chunks)

    rows = []
    for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        row = {
            "document_id": document_id,
            "content": chunk,
            "embedding": embedding,
            "chunk_index": idx,
            "metadata": {
                "filename": filename,
                "chunk_index": idx,
                "document_type": doc_meta.get("document_type", ""),
                "parties": doc_meta.get("parties", []),
            },
        }
        if user_id:
            row["user_id"] = user_id
        rows.append(row)

    supabase.table("chunks").insert(rows).execute()

    try:
        supabase.table("documents").update({"status": "completed"}).eq("id", document_id).execute()
    except Exception:
        logger.warning("Could not update document status for %s", document_id)

    logger.info("Upload '%s' ingested successfully (%d chunks)", filename, document_id)
    return len(chunks)
